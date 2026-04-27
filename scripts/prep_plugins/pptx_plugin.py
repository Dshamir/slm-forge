"""
PPTX plugin — per-slide text extraction via python-pptx.

Preserves behavior from prep-publications.py::emit_chunks_pptx exactly.
"""
from __future__ import annotations

import sys
from pathlib import Path
from typing import Iterator

from .orchestration_helpers import clean_text


class _PptxPlugin:
    extensions = (".pptx",)
    source_format = "pptx"
    requires = ("python-pptx",)
    system_deps = ()
    default_on = True
    disable_env = "FORGE_DISABLE_PPTX"

    def iter_chunks(self, path: Path, section: str, base_id: str, options: dict) -> Iterator[dict]:
        try:
            from pptx import Presentation
        except ImportError:
            sys.stderr.write(f"  [python-pptx not installed; skip {path}]\n")
            return
        try:
            prs = Presentation(str(path))
        except Exception as e:
            sys.stderr.write(f"  [{path}: {type(e).__name__}: {e}]\n")
            return
        for i, slide in enumerate(prs.slides):
            pieces = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs)
                    if t.strip():
                        pieces.append(t)
            if not pieces:
                continue
            text = clean_text("\n".join(pieces))
            if not text:
                continue
            yield {
                "id": f"{base_id}-s{i:03d}",
                "text": text,
                "format": "pretrain",
                "metadata": {
                    "source_file": str(path),
                    "source_format": "pptx",
                    "section": section,
                    "doc_title": path.stem,
                    "chunk_type": "slide",
                    "chunk_idx": i,
                    "char_count": len(text),
                },
            }


PLUGIN = _PptxPlugin()
